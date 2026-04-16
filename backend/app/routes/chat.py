import base64
import json
import os
import re
import tempfile
from typing import Annotated, Optional
from fastapi import APIRouter, Depends, HTTPException, File, Form, UploadFile
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sse_starlette.sse import EventSourceResponse
import httpx
from app.database import get_db
from app.models import User, Conversation, Message
from app.schemas import ChatRequest
from app.deps import get_auth_context, AuthContext
from app.llm_client import LLMClient
from app.field_types import build_prompt_field_types_compact
from app.context_compact import ContextCompactor
from app.json_utils import loads_if_str


async def _get_tenant_llm_config(db: AsyncSession, tenant_id: int) -> dict | None:
    from app.harness.llm_resolver import resolve_llm_config
    resolved = await resolve_llm_config(db, tenant_id, purpose="builder")
    if not resolved:
        return None
    return {"api_key": resolved.api_key, "base_url": resolved.base_url, "model": resolved.model, "max_tokens": resolved.max_tokens}


async def _get_conversation_llm_config(db: AsyncSession, conversation: Conversation) -> dict | None:
    from app.harness.llm_resolver import resolve_llm_config
    resolved = await resolve_llm_config(db, conversation.tenant_id, purpose="builder", selected_config_id=conversation.selected_llm_config_id)
    if not resolved:
        return None
    return {"api_key": resolved.api_key, "base_url": resolved.base_url, "model": resolved.model, "max_tokens": resolved.max_tokens}


async def _stream_with_tenant_llm(cfg: dict | None, messages: list, max_retries: int = 2):
    """使用租户 LLM 配置流式调用，yield OpenAI 格式 JSON 字符串。网络错误自动重试。"""
    import logging as _logging
    _logger = _logging.getLogger(__name__)

    if cfg is None:
        raise ValueError("未配置可用的 LLM 模型，请在环境管理 → 模型配置中添加并启用模型")

    llm = LLMClient(api_key=cfg["api_key"], base_url=cfg["base_url"], model=cfg["model"])
    _logger.info("LLM stream call: model=%s base_url=%s", cfg["model"], cfg["base_url"])

    import asyncio
    last_err = None
    for attempt in range(max_retries + 1):
        if attempt > 0:
            await asyncio.sleep(1)
            _logger.info("LLM stream retry %d/%d", attempt, max_retries)
        try:
            async for chunk in llm.chat_completion_stream(messages, max_tokens=cfg.get("max_tokens", 8192)):
                yield chunk
            return
        except (httpx.ConnectError, httpx.ConnectTimeout, httpx.ReadTimeout) as e:
            last_err = e
            _logger.warning("LLM stream attempt %d failed: %r (type=%s)", attempt + 1, e, type(e).__name__)
            continue
        except Exception:
            raise

    if last_err:
        raise last_err

router = APIRouter(prefix="/chat", tags=["聊天"])


def _strip_hidden_blocks(text: str) -> str:
    cleaned = str(text or "")
    cleaned = re.sub(r"<think>[\s\S]*?</think>", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"<think>[\s\S]*$", "", cleaned, flags=re.IGNORECASE)
    return cleaned.strip()


async def _parse_uploaded_document(file: UploadFile) -> str:
    """尽量从常见文档中提取可读文本，提取失败时返回空串。"""
    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()
    content_type = (file.content_type or "").lower()
    content = await file.read()

    text_like_exts = {
        ".md", ".markdown", ".txt", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm"
    }
    if ext in text_like_exts or content_type.startswith("text/"):
        return content.decode("utf-8", errors="ignore")

    suffix = ext or ".tmp"
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        if ext == ".pdf":
            try:
                import pdfplumber
                parts: list[str] = []
                with pdfplumber.open(tmp_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            parts.append(text)
                return "\n".join(parts)
            except Exception:
                return ""

        if ext in {".docx", ".doc"}:
            try:
                from docx import Document
                doc = Document(tmp_path)
                parts: list[str] = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            parts.append(row_text)
                return "\n".join(parts)
            except Exception:
                return ""

        if ext in {".xlsx", ".xls"}:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(tmp_path, data_only=True)
                parts: list[str] = []
                for ws in wb.worksheets:
                    parts.append(f"## Sheet: {ws.title}")
                    for row in ws.iter_rows(values_only=True):
                        values = [str(v).strip() for v in row if v is not None and str(v).strip()]
                        if values:
                            parts.append(" | ".join(values))
                return "\n".join(parts)
            except Exception:
                return ""

        if ext in {".pptx", ".ppt"}:
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                parts: list[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    slide_texts: list[str] = []
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            slide_texts.append(text.strip())
                    if slide_texts:
                        parts.append(f"## Slide {idx}")
                        parts.extend(slide_texts)
                return "\n".join(parts)
            except Exception:
                return ""

        return ""
    finally:
        os.unlink(tmp_path)

BUILDER_SYSTEM_PROMPT = """你是 aPaaS Builder AI，得帆云低代码平台的智能搭建助手。你的任务是通过多轮对话帮用户理清应用需求，然后生成结构化的应用配置。

## 平台知识

得帆云 aPaaS 是企业级低代码平台，核心架构：平台 → 租户 → 应用 → 菜单 → 表单 → 数据模型。

**7种菜单页面类型**：表单页面（最常用，自动生成数据模型）、模型页面（复用已有模型）、聚合表单（跨表只读汇总）、分析页面（可视化图表）、外部链接、自开发页面、引用功能菜单。

**表单构建四步流程**：表单设计 → 列表设计 → 流程设计 → 页面设置。

**组件分类**：
- 业务组件：单行/多行输入、数字、日期时间、人员/部门选择、单选/多选/下拉、金额、附件、富文本、地址、定位、超链接、开关等
- 高级组件：单据号（自增编号）、数据选择、数据单选、关联表单、数据统计、他表字段、子表、虚拟字段
- 子表适用于明细行（订单明细、配件清单等），会生成独立的子表数据模型

**数据字典**：用于下拉选择类字段的选项管理，支持应用级和租户级共享。

**权限体系**：
- 功能权限：查看、新增、编辑、删除、导入、导出等操作控制
- 数据权限：控制可见数据范围（本人/本部门/全部等）
- 数据范围类型：SELF(本人)、CURRENT_USER_DEPT(本部门)、CURRENT_USER_DEPT_LOW_LEVEL(本部门及下级)、ALL(全部)

**表结构设计规范**：
- 平台自动维护以下系统字段，**设计表结构时绝对不要添加**：id、创建时间、更新时间、创建人、更新人（以及任何同义变体如 create_time、update_by 等）
- 只设计真正的业务字段

**角色设计规范**：
- "员工"、"全体员工"、"所有员工" 等通用性角色**不需要创建**，平台内置"全体成员"概念，在权限配置中用 role="all" 表达
- "直属上级"、"部门经理"、"上级领导" 等层级关系角色**不需要创建**，平台直接读取组织架构，在审批流程节点中配置即可

**重要限制**：
- 应用编码和菜单编码创建后不可修改
- 修改表单后需重新发布才在前台生效
- code 只能用英文字母、数字和下划线，不能用中文，必须以字母开头，不能是数据库保留字

## 对话流程（三阶段）

### 阶段一：需求分析（当用户给的是模糊业务需求时）
如果用户描述的是业务场景而非详细功能设计（如"做个售后管理系统"、"客户报修后派工程师上门"），你需要先做需求拆解：

1. **识别业务实体**：从业务描述中提取核心对象（如：工单、客户、工程师）
2. **推导字段和字典**：为每个实体推导关键字段及其类型、枚举选项（如：工单状态 = 待派工/维修中/已完成）
3. **推导表间关系**：实体之间如何关联（如：工单关联客户、工单关联工程师）

**输出格式**：用 2-3 轮对话完成拆解，每轮给出你的分析并请用户确认/补充：
```
基于你的描述，我分析出以下核心功能：

**核心模块：**
- 工单管理
- 客户档案
- 工程师信息

**需要你确认：**
1. 派工方式是手动选人还是需要智能派单？
2. 工程师上门需要拍照记录吗？
3. 还有哪些角色？客服、主管？
```

### 阶段二：需求细化（对话引导）
通过提问理清需求细节：数据模型、字段、角色、字典、权限。

### 阶段三：配置生成
需求明确后，生成完整的配置JSON。

## 重要规则
- 用中文回复，使用markdown格式
- 主动引导用户补充信息（特别是：哪些字段需要枚举选项？哪些表单之间有关联？）
- 模型之间有关联时，使用"数据单选"字段类型并指定ref
- 需要枚举选项的字段，**必须**先在dicts中定义字典，再在字段中用dict引用
- 有明细行的表单（如订单明细），使用"子表"字段类型
- 所有code字段只用英文、数字、下划线，以字母开头
- code必须避免数据库保留字，如name/status/type/order/date/number/code/description/title/content/note/remark/contact/price/total/quantity/company/customer/product/service/region等。建议加业务前缀，如customer_name、order_status
- **【禁止】** 在任何表单的 fields 中出现 id、创建时间、更新时间、创建人、更新人及其英文变体（create_time、update_time、created_by、updated_by、creator 等），这些由平台自动维护
- **【禁止】** 创建"员工"、"全体员工"等通用性角色，需要表达全员时直接在 permissions 中使用 role="all"
- **【禁止】** 创建"直属上级"、"部门经理"等层级角色，层级关系由平台组织架构自动处理
- **【必须】** 每个表单的 permissions 中，**默认包含一条 role="all" 的全体人员规则**，再叠加其他角色的差异化权限

## 生成配置
当需求明确后，在回复最后附上配置JSON，用 ```json 代码块包裹：

```json
{"type":"preview","data":{"appName":"应用名","roles":[{"code":"role_admin","name":"管理员"}],"dicts":[{"name":"字典名","code":"dict_code","options":[{"name":"选项名","code":"opt_code"}]}],"models":[{"name":"表单名","code":"model_code","fields":[{"name":"字段名","type":"单行输入","icon":"T","required":true,"code":"field_code","dict":"dict_code","ref":{"model":"other_model_code","field":"display_field_code"},"sub_code":"sub_model_code","sub_fields":[{"name":"子字段名","type":"单行输入","icon":"T","code":"sub_field_code"}]}]}],"permissions":[{"form":"表单名","rules":[{"role":"all","op":"all","data":"ALL"}]}]}}
```

## 字段类型（type可选值及icon）
""" + build_prompt_field_types_compact() + """

## 字段属性说明
- code: 字段编码，英文+下划线，如 asset_name
- dict: 下拉单选/多选时，填写dicts中定义的字典code
- ref: 数据单选时，填写 {"model":"关联模型code","field":"要显示的字段code"}
- sub_code: 子表时，填写子表模型code
- sub_fields: 子表内的字段数组（格式同普通字段）

## 权限配置说明
- role: 角色code，"all" 表示全体人员（平台内置，**无需在 roles 中定义**）
- op: 操作类型，"all"=全部操作, "add"=新增, "edit"=编辑, "delete"=删除
- data: 数据范围，"ALL"=全部, "SELF"=本人, "dept"=本部门
- **每个表单权限必须包含全体人员兜底规则**，示例：`{"role":"all","op":"all","data":"ALL"}`，再按需叠加角色限制"""

ASSISTANT_SYSTEM_PROMPT = """你是 aPaaS 辅助开发智能体，帮助用户完善已有应用的配置，包括：创建审批流程、配置业务规则、调整表单组件、配置数据权限。用中文回复，使用markdown格式。"""

DEVELOPER_SYSTEM_PROMPT = """你是 aPaaS 复杂开发智能体，帮助用户进行二次开发，包括：自定义Vue组件、后端接口集成、Groovy脚本编写。用中文回复，使用markdown格式。"""

REQUIREMENTS_SYSTEM_PROMPT = """你是一个懂业务、懂系统的应用搭建顾问。用户可能是业务人员，也可能是 IT，他们都有一个共同点：脑子里有想法，但还没整理清楚。你的工作是通过自然对话帮他们把想法变成一个可运行的应用。

## 你需要搞清楚的事（内部目标，不要展示给用户）

- 这个系统解决什么业务问题
- 谁会用，不同人看到的数据范围是否不同
- 要管理哪些业务对象（每个对象会变成一张表单）
- 每个对象有哪些关键信息要记录
- 哪些字段是选项类的（比如状态、类型），选项有哪些
- 不同角色能做什么操作，能看多大范围的数据

平台能力边界（用户提到时如实说明，不要追问细节）：
- 消息通知、定时提醒、外部系统对接：平台暂不支持，跳过
- 审批流程的状态（待审批/已通过等）：平台内置，不需要用户定义
- "全体员工""直属上级"这类泛化角色：平台内置，不需单独创建

## 对话方式

每次只问一个问题，根据用户的回答自然推进，像顾问聊需求一样。对业务人员用业务语言，对 IT 用户可以稍微技术一点，但都不要暴露"枚举值""数据模型"这类系统术语。

示例：
- 不说"枚举值有哪些" → 说"状态有几种？比如待处理、进行中、已完成？"
- 不说"数据模型" → 说"要管理哪些信息，比如客户、订单、设备这类"
- 不说"权限配置" → 说"谁能看到所有数据，谁只能看自己的"

## 信息够了就生成设计文档

当你判断核心信息已经清晰（知道做什么、管理哪些对象、谁用），不要继续追问，直接输出设计文档让用户确认。没问到的内容根据业务场景合理推断，推断的内容在文档里体现出来让用户看到。

输出格式（结构化 Markdown）：

---

## 功能设计文档

### 应用概述
- **应用名称**：xxx
- **应用编码**：XXX（英文大写缩写）
- **描述**：xxx

### 角色清单
| 角色编码 | 角色名称 | 数据范围 |
|---------|---------|---------|
| xxx | xxx | 仅本人/本部门/全公司 |

### 数据字典
**字典名（dict_code）**：选项1、选项2、选项3
（每个字典一行）

### 数据模型
**表名（table_code）**[主表/子表]
描述：xxx
字段：字段1（类型）、字段2（类型）、...

### 权限设计
| 数据表 | 角色 | 操作权限 | 数据范围 |
|-------|------|---------|---------|
| xxx | xxx | 新增/编辑/查看 | 仅本人/本部门/全公司 |

---

输出文档后最后一行加：`<!-- DESIGN_COMPLETE -->`

用户确认后（说"可以""没问题""开始生成"等），回复"好的，开始生成应用配置。"并在最后加：`<!-- TRIGGER_BUILD -->`

**注意：只输出 Markdown，不输出 JSON；不要说"我无法创建应用"。**"""

SYSTEM_PROMPTS = {
    "builder": BUILDER_SYSTEM_PROMPT,
    "requirements": REQUIREMENTS_SYSTEM_PROMPT,
    "assistant": ASSISTANT_SYSTEM_PROMPT,
    "developer": DEVELOPER_SYSTEM_PROMPT,
}


def _build_incremental_config_prompt(current_config: dict | None) -> str:
    if not current_config:
        return ""
    data = current_config.get("data", current_config) if isinstance(current_config, dict) else {}
    app_name = data.get("appName") or data.get("app_name") or "当前应用"
    roles = data.get("roles") or []
    dicts = data.get("dicts") or []
    models = data.get("models") or []
    forms = data.get("forms") or []
    permissions = data.get("permissions") or []
    return (
        "你正在基于已有应用配置做增量调整，不要重新解析整份设计文档，不要重新输出“请确认以下信息”这类全文确认内容。"
        f"\n当前应用：{app_name}"
        f"\n当前配置概况：{len(roles)} 个角色，{len(dicts)} 个字典，{len(models)} 个模型，{len(forms)} 个表单，{len(permissions)} 个权限组。"
        "\n请只围绕用户本次补充的内容进行修改、补充、解释或给出增量配置建议。"
        "\n如果用户是在修改配置，请不要返回完整 preview JSON。"
        "\n必须只输出一个 ```json 代码块，格式为："
        '\n{"type":"patch","actions":[...]}'
        "\n只允许使用这些 patch op："
        "\nadd_dict, update_dict, remove_dict,"
        "\nadd_field, update_field, remove_field,"
        "\nadd_model, remove_model,"
        "\nadd_role, remove_role,"
        "\nadd_workflow, update_workflow, remove_workflow,"
        "\nadd_permission, update_permission, remove_permission, set_permissions。"
        "\n禁止输出完整 preview，禁止整份重写当前 JSON。"
        "\n如果只是解释，不需要修改配置，则不要输出 JSON。"
    )


def _extract_config_from_response(content: str):
    """从 LLM 响应中提取 preview 或 patch JSON，返回 (type, data)。"""
    import re as _re
    for m in _re.finditer(r'```json\s*([\s\S]*?)\s*```', content):
        try:
            obj = json.loads(m.group(1))
            if isinstance(obj, dict) and obj.get("type") in ("preview", "patch"):
                return obj["type"], obj
        except Exception:
            pass
    return None, None


def _apply_patch_to_config(config: dict, patch: dict) -> dict:
    """将 patch actions 应用到本地 config，返回更新后的 config。"""
    import copy
    result = copy.deepcopy(config)
    data = result.get("data", result) if isinstance(result, dict) else result

    for action in patch.get("actions", []):
        op = action.get("op", "")
        try:
            if op == "add_role":
                data.setdefault("roles", []).append({"name": action["name"], "code": action["code"]})
            elif op == "remove_role":
                data["roles"] = [r for r in data.get("roles", []) if r.get("code") != action.get("code")]
            elif op == "add_dict":
                data.setdefault("dicts", []).append({
                    "name": action["name"], "code": action["code"],
                    "options": action.get("options", [])
                })
            elif op == "update_dict":
                for d in data.get("dicts", []):
                    if d.get("code") == action.get("code"):
                        if "name" in action: d["name"] = action["name"]
                        if "options" in action: d["options"] = action["options"]
            elif op == "remove_dict":
                data["dicts"] = [d for d in data.get("dicts", []) if d.get("code") != action.get("code")]
            elif op == "add_model":
                data.setdefault("models", []).append({
                    "name": action["name"], "code": action["code"],
                    "fields": action.get("fields", [])
                })
            elif op == "remove_model":
                data["models"] = [m for m in data.get("models", []) if m.get("code") != action.get("code")]
            elif op == "add_field":
                for model in data.get("models", []):
                    if model.get("code") == action.get("model"):
                        model.setdefault("fields", []).append(action.get("field", {}))
            elif op == "update_field":
                for model in data.get("models", []):
                    if model.get("code") == action.get("model"):
                        for field in model.get("fields", []):
                            if field.get("code") == action.get("field"):
                                field.update(action.get("changes", {}))
            elif op == "remove_field":
                for model in data.get("models", []):
                    if model.get("code") == action.get("model"):
                        model["fields"] = [f for f in model.get("fields", []) if f.get("code") != action.get("field")]
            elif op in ("add_permission", "set_permissions"):
                perms = data.setdefault("permissions", [])
                form_name = action.get("form")
                existing = next((p for p in perms if p.get("form") == form_name), None)
                if existing:
                    existing["rules"] = action.get("rules", existing.get("rules", []))
                else:
                    perms.append({"form": form_name, "rules": action.get("rules", [])})
            elif op == "update_permission":
                for p in data.get("permissions", []):
                    if p.get("form") == action.get("form"):
                        p["rules"] = action.get("rules", p.get("rules", []))
            elif op == "remove_permission":
                data["permissions"] = [p for p in data.get("permissions", []) if p.get("form") != action.get("form")]
        except Exception:
            pass  # 单个 action 失败不影响其他
    return result


def _resolve_effective_config(server_config, client_config_raw) -> dict | None:
    """优先用服务端 config；没有时兼容旧版前端传参。"""
    if server_config is not None:
        return server_config
    if not client_config_raw:
        return None
    try:
        return loads_if_str(client_config_raw)
    except Exception:
        return None


def _build_phase_prompt(agent_type: str, effective_config) -> str:
    """根据 agent_type 和当前 config 决定使用哪个 system prompt。"""
    # requirements agent 只做需求对话，不做增量 config 修改
    if agent_type == "requirements":
        return REQUIREMENTS_SYSTEM_PROMPT

    if agent_type != "builder":
        base = SYSTEM_PROMPTS.get(agent_type, BUILDER_SYSTEM_PROMPT)
        incr = _build_incremental_config_prompt(effective_config)
        return base + ("\n\n" + incr if incr else "")

    # builder：有 config → refining，没有 → gathering
    if effective_config:
        incr = _build_incremental_config_prompt(effective_config)
        return BUILDER_SYSTEM_PROMPT + ("\n\n" + incr if incr else "")
    else:
        return REQUIREMENTS_SYSTEM_PROMPT


@router.post("/send")
async def send_message(
    data: ChatRequest,
    ctx: Annotated[AuthContext, Depends(get_auth_context)],
    db: Annotated[AsyncSession, Depends(get_db)]
):
    # 验证对话是否存在且属于当前用户（租户隔离）
    result = await db.execute(
        select(Conversation).where(
            Conversation.id == data.conversation_id,
            Conversation.user_id == ctx.user.id,
            Conversation.tenant_id == ctx.tenant_id
        )
    )
    conversation = result.scalar_one_or_none()

    if not conversation:
        raise HTTPException(status_code=404, detail="对话不存在")

    # 保存用户消息
    user_message = Message(
        conversation_id=conversation.id,
        role="user",
        content=data.message
    )
    db.add(user_message)
    await db.commit()

    # 自动更新对话标题（首条用户消息时）
    if conversation.title in ("新对话", "需求分析", "智能开发") or conversation.title.startswith("新对话"):
        short_title = data.message.strip().replace("\n", " ")[:30]
        if short_title:
            conversation.title = short_title
            await db.commit()

    # 获取历史消息
    result = await db.execute(
        select(Message)
        .where(Message.conversation_id == conversation.id)
        .order_by(Message.created_at)
    )
    messages = result.scalars().all()

    # 服务端 config 优先，兼容旧版前端传参
    effective_config = _resolve_effective_config(
        conversation.current_config,
        getattr(data, 'current_config', None)
    )

    # 根据 phase 选 system prompt
    system_prompt = _build_phase_prompt(conversation.agent_type, effective_config)
    llm_messages = [{"role": "system", "content": system_prompt}]

    # 预取租户 LLM 配置（必须在 compactor 之前）
    llm_cfg = await _get_conversation_llm_config(db, conversation)

    # 智能上下文压缩
    history_msgs = [{"role": msg.role, "content": msg.content} for msg in messages]
    compactor = ContextCompactor(llm_cfg)
    compacted, new_summary = await compactor.compact_with_summary(
        history_msgs,
        mode="chat",
        existing_summary=getattr(conversation, 'context_summary', None),
    )
    if new_summary and new_summary != getattr(conversation, 'context_summary', None):
        try:
            conversation.context_summary = new_summary
            await db.commit()
        except Exception:
            pass
    llm_messages.extend(compacted)

    # 流式响应
    async def event_generator():
        assistant_content = ""
        thinking_sent = False

        try:
            async for chunk in _stream_with_tenant_llm(llm_cfg, llm_messages):
                chunk_data = json.loads(chunk)
                if "choices" in chunk_data and len(chunk_data["choices"]) > 0:
                    delta = chunk_data["choices"][0].get("delta", {})
                    content = delta.get("content") or ""

                    has_reasoning = bool(delta.get("reasoning_details") or delta.get("reasoning_content"))
                    if has_reasoning and not thinking_sent:
                        thinking_sent = True
                        yield {"event": "thinking", "data": json.dumps({"type": "thinking", "data": "思考中..."})}

                    if content:
                        assistant_content += content
                        visible_content = _strip_hidden_blocks(content)
                        if visible_content:
                            yield {"event": "message", "data": json.dumps({"type": "message", "data": visible_content})}

            # 保存助手消息（strip think 块后存库）
            assistant_content = _strip_hidden_blocks(assistant_content)
            db.add(Message(conversation_id=conversation.id, role="assistant", content=assistant_content))
            await db.commit()

            # 提取响应中的 config（preview 或 patch），存服务端并通知前端
            cfg_type, cfg_data = _extract_config_from_response(assistant_content)
            if cfg_type == "preview" and cfg_data:
                try:
                    conversation.current_config = cfg_data
                    conversation.phase = "refining"
                    await db.commit()
                except Exception:
                    pass
                yield {"event": "config", "data": json.dumps(
                    {"type": "preview", "data": cfg_data.get("data", cfg_data)},
                    ensure_ascii=False
                )}
            elif cfg_type == "patch" and cfg_data and effective_config:
                try:
                    new_config = _apply_patch_to_config(effective_config, cfg_data)
                    conversation.current_config = new_config
                    await db.commit()
                    yield {"event": "config", "data": json.dumps(
                        {"type": "preview", "data": new_config.get("data", new_config)},
                        ensure_ascii=False
                    )}
                except Exception:
                    pass

            yield {"event": "done", "data": json.dumps({"type": "done", "data": "completed"})}

        except Exception as e:
            yield {"event": "error", "data": json.dumps({"type": "error", "data": str(e)})}

    return EventSourceResponse(event_generator())


@router.post("/send-with-file")
async def send_message_with_file(
    ctx: Annotated[AuthContext, Depends(get_auth_context)],
    db: Annotated[AsyncSession, Depends(get_db)],
    conversation_id: int = Form(...),
    message: str = Form(default=""),
    file: Optional[UploadFile] = File(default=None),
    current_config: str = Form(default=""),
):
    _IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp"}
    _MIME_MAP = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".gif": "image/gif",
        ".webp": "image/webp",
    }

    result = await db.execute(
        select(Conversation).where(
            Conversation.id == conversation_id,
            Conversation.user_id == ctx.user.id,
            Conversation.tenant_id == ctx.tenant_id
        )
    )
    conversation = result.scalar_one_or_none()
    if not conversation:
        raise HTTPException(status_code=404, detail="对话不存在")

    image_data_url = ""
    file_content = ""
    file_name = ""
    route_to_pipeline = False  # 是否把文件路由到 doc_pipeline

    if file and file.filename:
        file_name = file.filename
        ext = os.path.splitext(file_name)[1].lower()
        if ext in _IMAGE_EXTS:
            raw = await file.read()
            mime = _MIME_MAP.get(ext, "image/png")
            b64 = base64.b64encode(raw).decode()
            image_data_url = f"data:{mime};base64,{b64}"
        else:
            file_content = await _parse_uploaded_document(file)
            # 对文档文件做标准度检测，决定走 pipeline 还是 chat 上下文
            if file_content and conversation.agent_type in {"builder", "requirements"}:
                try:
                    from app.doc_standard_detector import detect as _detect_doc
                    doc_score = _detect_doc(file_content).get("score", 0)
                    if doc_score >= 60:
                        route_to_pipeline = True
                except Exception:
                    pass

    if not message.strip() and not image_data_url and not file_name:
        raise HTTPException(status_code=400, detail="消息不能为空")

    db_content = message.strip()
    if file_content and not route_to_pipeline:
        db_content = f"{db_content}\n\n[上传文件：{file_name}]\n\n{file_content}".strip()
    elif file_name:
        tag = "上传截图" if image_data_url else "上传文件"
        db_content = f"{db_content}\n\n[{tag}：{file_name}]".strip()

    db.add(Message(conversation_id=conversation.id, role="user", content=db_content))
    await db.commit()

    if conversation.title in ("新对话", "需求分析", "智能开发") or conversation.title.startswith("新对话"):
        short_title = (message.strip() or f"文件：{file_name}").replace("\n", " ")[:30]
        if short_title:
            conversation.title = short_title
            await db.commit()

    llm_cfg = await _get_conversation_llm_config(db, conversation)

    # ── 走 doc_pipeline（文档标准度 >= 60）───────────────────────
    if route_to_pipeline:
        import asyncio as _asyncio

        async def pipeline_event_generator():
            try:
                from app.doc_pipeline import parse_document
                progress_queue = _asyncio.Queue()

                async def _on_progress(msg: str, *, batch=None):
                    await progress_queue.put((msg, batch))

                parse_task = _asyncio.create_task(
                    parse_document(file_content, llm_cfg=llm_cfg, on_progress=_on_progress)
                )

                yield {"event": "message", "data": json.dumps(
                    {"type": "message", "data": f"检测到设计文档（{file_name}），正在解析...\n\n"},
                    ensure_ascii=False
                )}

                while not parse_task.done():
                    try:
                        msg, batch = await _asyncio.wait_for(progress_queue.get(), timeout=0.3)
                        event_data: dict = {"message": msg}
                        if batch is not None:
                            # 从 msg tag 中推断模块类型，附带局部数据让前端实时渲染
                            for key in ("roles", "dicts", "models", "forms", "permissions"):
                                if f"[{key}]" in msg:
                                    event_data["module"] = key
                                    event_data["data"] = batch
                                    break
                        yield {"event": "progress", "data": json.dumps(event_data, ensure_ascii=False)}
                    except _asyncio.TimeoutError:
                        continue

                parse_result = await parse_task
                # 清空任务完成后队列中残留的 progress 项
                while not progress_queue.empty():
                    try:
                        msg, batch = progress_queue.get_nowait()
                        event_data = {"message": msg}
                        if batch is not None:
                            for key in ("roles", "dicts", "models", "forms", "permissions"):
                                if f"[{key}]" in msg:
                                    event_data["module"] = key
                                    event_data["data"] = batch
                                    break
                        yield {"event": "progress", "data": json.dumps(event_data, ensure_ascii=False)}
                    except Exception:
                        break
                config_data = parse_result if isinstance(parse_result, dict) else {}

                # 存服务端
                conversation.current_config = config_data
                conversation.phase = "refining"
                db.add(Message(
                    conversation_id=conversation.id,
                    role="assistant",
                    content=f"[文档解析完成：{file_name}]"
                ))
                await db.commit()

                preview_data = config_data.get("data", config_data)
                yield {"event": "config", "data": json.dumps(
                    {"type": "preview", "data": preview_data}, ensure_ascii=False
                )}
                yield {"event": "done", "data": json.dumps({"type": "done", "data": "completed"})}

            except Exception as e:
                yield {"event": "error", "data": json.dumps({"type": "error", "data": str(e)})}

        return EventSourceResponse(pipeline_event_generator())

    # 判断是否走"文档静默生成"路径
    effective_config = _resolve_effective_config(conversation.current_config, current_config)
    _in_requirements_phase = (
        conversation.agent_type == "requirements" or
        (conversation.agent_type == "builder" and not effective_config)
    )
    _doc_is_substantial = len(file_content) > 1500

    # ── 走文档静默生成（文档够详细，不需要对话确认）────────────────────
    if file_content and _in_requirements_phase and _doc_is_substantial:
        import asyncio as _asyncio

        async def doc_silent_generator():
            try:
                # 1. 告诉用户正在生成，不把 AI 分析文字流式输出
                yield {"event": "message", "data": json.dumps(
                    {"type": "message", "data": f"正在解析文档「{file_name}」，生成应用配置...\n"},
                    ensure_ascii=False
                )}
                # 每隔一段发一个 progress 保持连接
                keep_alive_task = _asyncio.create_task(_asyncio.sleep(0))  # placeholder

                silent_prompt = BUILDER_SYSTEM_PROMPT + (
                    "\n\n【文档直出模式】用户上传了一份完整的需求/设计文档，文档内容已附在消息末尾。"
                    "请直接阅读文档，一次性输出完整的应用配置 JSON（type=preview）。"
                    "不要输出分析过程、不要追问用户、不要输出 Markdown 设计文档。"
                    "文档中未明确说明的细节按业务场景合理补全。"
                )
                silent_messages = [
                    {"role": "system", "content": silent_prompt},
                    {"role": "user", "content": db_content},
                ]

                # 2. 累积完整响应（流式接收但不向前端输出文字）
                full_content = ""
                async for chunk in _stream_with_tenant_llm(llm_cfg, silent_messages):
                    chunk_data = json.loads(chunk)
                    if "choices" in chunk_data and chunk_data["choices"]:
                        delta = chunk_data["choices"][0].get("delta", {})
                        full_content += delta.get("content") or ""
                    # 每积累 500 字发一个 progress 保持连接活跃
                    if len(full_content) % 500 < 10:
                        yield {"event": "progress", "data": json.dumps(
                            {"message": f"正在生成配置... ({len(full_content)} 字)"}, ensure_ascii=False
                        )}

                full_content = _strip_hidden_blocks(full_content)

                # 3. 提取 JSON 配置
                cfg_type, cfg_data = _extract_config_from_response(full_content)
                if cfg_type == "preview" and cfg_data:
                    conversation.current_config = cfg_data
                    conversation.phase = "refining"
                    db.add(Message(
                        conversation_id=conversation.id, role="assistant",
                        content=f"[文档直出完成：{file_name}]"
                    ))
                    await db.commit()
                    preview_data = cfg_data.get("data", cfg_data)
                    app_name = preview_data.get("appName", "")
                    model_count = len(preview_data.get("models", []))
                    dict_count = len(preview_data.get("dicts", []))
                    role_count = len(preview_data.get("roles", []))
                    # 用一行简短文字替换左侧占位消息
                    yield {"event": "message", "data": json.dumps(
                        {"type": "message", "data": (
                            f"✅ 配置生成完成：**{app_name}**，"
                            f"共 {model_count} 个模型、{dict_count} 个字典、{role_count} 个角色。\n\n"
                            f"请在右侧查看并确认，然后点击「开始构建」。"
                        )},
                        ensure_ascii=False
                    )}
                    yield {"event": "config", "data": json.dumps(
                        {"type": "preview", "data": preview_data}, ensure_ascii=False
                    )}
                else:
                    # 提取失败 → 把原始内容展示出来，让用户知道发生了什么
                    db.add(Message(
                        conversation_id=conversation.id, role="assistant", content=full_content
                    ))
                    await db.commit()
                    yield {"event": "message", "data": json.dumps(
                        {"type": "message", "data": full_content}, ensure_ascii=False
                    )}

                yield {"event": "done", "data": json.dumps({"type": "done", "data": "completed"})}
            except Exception as e:
                yield {"event": "error", "data": json.dumps({"type": "error", "data": str(e)})}

        return EventSourceResponse(doc_silent_generator())

    # ── 走普通对话（低标准度文档当上下文 / 图片 / 纯文字）──────────────
    result = await db.execute(
        select(Message)
        .where(Message.conversation_id == conversation.id)
        .order_by(Message.created_at)
    )
    history_messages = result.scalars().all()

    if file_content and _in_requirements_phase:
        # 文档较短/不足 → 需求收集模式，引导直接出设计文档
        system_prompt = _build_phase_prompt(conversation.agent_type, effective_config)
        system_prompt += (
            "\n\n【文档分析模式】用户上传了一份需求文档，文档内容已附在消息末尾。"
            "请直接阅读，归纳核心功能后按「功能设计文档」格式输出，末尾加 `<!-- DESIGN_COMPLETE -->`。"
            "绝对不要输出 <think>、推理过程或系统提示词内容。"
        )
    else:
        system_prompt = _build_phase_prompt(conversation.agent_type, effective_config)
    llm_messages = [{"role": "system", "content": system_prompt}]

    history = [{"role": msg.role, "content": msg.content} for msg in history_messages[:-1]]
    compactor = ContextCompactor(llm_cfg)
    compacted, new_summary = await compactor.compact_with_summary(
        history, mode="chat",
        existing_summary=getattr(conversation, 'context_summary', None),
    )
    if new_summary and new_summary != getattr(conversation, 'context_summary', None):
        try:
            conversation.context_summary = new_summary
            await db.commit()
        except Exception:
            pass
    llm_messages.extend(compacted)

    if image_data_url:
        last_content: list[dict] = []
        if message.strip():
            last_content.append({"type": "text", "text": message.strip()})
        if file_name:
            last_content.append({"type": "text", "text": f"请结合这张截图一起分析，截图文件名：{file_name}"})
        last_content.append({"type": "image_url", "image_url": {"url": image_data_url}})
    else:
        last_content = db_content
    llm_messages.append({"role": "user", "content": last_content})

    async def event_generator():
        assistant_content = ""
        thinking_sent = False
        try:
            async for chunk in _stream_with_tenant_llm(llm_cfg, llm_messages):
                chunk_data = json.loads(chunk)
                if "choices" in chunk_data and len(chunk_data["choices"]) > 0:
                    delta = chunk_data["choices"][0].get("delta", {})
                    content = delta.get("content") or ""
                    has_reasoning = bool(delta.get("reasoning_details") or delta.get("reasoning_content"))
                    if has_reasoning and not thinking_sent:
                        thinking_sent = True
                        yield {"event": "thinking", "data": json.dumps({"type": "thinking", "data": "思考中..."})}
                    if content:
                        assistant_content += content
                        visible_content = _strip_hidden_blocks(content)
                        if visible_content:
                            yield {"event": "message", "data": json.dumps({"type": "message", "data": visible_content})}

            assistant_content = _strip_hidden_blocks(assistant_content)
            db.add(Message(conversation_id=conversation.id, role="assistant", content=assistant_content))
            await db.commit()

            # 提取 config 存服务端并通知前端
            cfg_type, cfg_data = _extract_config_from_response(assistant_content)
            if cfg_type == "preview" and cfg_data:
                try:
                    conversation.current_config = cfg_data
                    conversation.phase = "refining"
                    await db.commit()
                except Exception:
                    pass
                yield {"event": "config", "data": json.dumps(
                    {"type": "preview", "data": cfg_data.get("data", cfg_data)}, ensure_ascii=False
                )}
            elif cfg_type == "patch" and cfg_data and effective_config:
                try:
                    new_config = _apply_patch_to_config(effective_config, cfg_data)
                    conversation.current_config = new_config
                    await db.commit()
                    yield {"event": "config", "data": json.dumps(
                        {"type": "preview", "data": new_config.get("data", new_config)}, ensure_ascii=False
                    )}
                except Exception:
                    pass

            yield {"event": "done", "data": json.dumps({"type": "done", "data": "completed"})}
        except Exception as e:
            yield {"event": "error", "data": json.dumps({"type": "error", "data": str(e)})}

    return EventSourceResponse(event_generator())


@router.post("/generate-config")
async def generate_config_phased(
    data: ChatRequest,
    ctx: Annotated[AuthContext, Depends(get_auth_context)],
    db: Annotated[AsyncSession, Depends(get_db)]
):
    """分阶段生成配置 (SSE)。

    前端在用户需求明确后调用此端点，替代让 LLM 一次性输出完整 JSON。
    分3阶段: 骨架 → 字典 → 模型，每阶段 yield 进度，前端实时更新预览。
    """
    from app.config_assembler import assemble_config_streaming

    # 验证对话
    result = await db.execute(
        select(Conversation).where(
            Conversation.id == data.conversation_id,
            Conversation.user_id == ctx.user.id,
            Conversation.tenant_id == ctx.tenant_id
        )
    )
    conversation = result.scalar_one_or_none()
    if not conversation:
        raise HTTPException(status_code=404, detail="对话不存在")

    # 收集历史上下文（摘要，不是全部消息）
    result = await db.execute(
        select(Message)
        .where(Message.conversation_id == conversation.id)
        .order_by(Message.created_at)
        .limit(20)
    )
    messages_list = result.scalars().all()
    context = "\n".join(
        f"{'用户' if m.role == 'user' else 'AI'}: {m.content[:500]}"
        for m in messages_list
        if m.role in ("user", "assistant")
    )

    llm_cfg = await _get_conversation_llm_config(db, conversation)

    async def event_generator():
        try:
            complete_config = None
            async for event in assemble_config_streaming(data.message, context, llm_cfg=llm_cfg):
                yield {
                    "event": "progress",
                    "data": json.dumps(event, ensure_ascii=False)
                }
                if event.get("phase") == "complete":
                    complete_config = event.get("data")

            # 保存完整配置为 system 消息，便于历史恢复
            if complete_config:
                config_msg = Message(
                    conversation_id=conversation.id,
                    role="system",
                    content=f"```json\n{json.dumps({'type': 'preview', 'data': complete_config}, ensure_ascii=False)}\n```"
                )
                db.add(config_msg)

                # 保存助手摘要消息
                summary = (
                    f"配置生成完成！共 {len(complete_config.get('models', []))} 个模型、"
                    f"{len(complete_config.get('dicts', []))} 个字典、"
                    f"{len(complete_config.get('roles', []))} 个角色。\n\n"
                    f"右侧预览已更新，你可以：\n"
                    f"1. 点击字典/角色上的编辑按钮直接修改\n"
                    f"2. 告诉我需要调整的内容\n"
                    f"3. 确认无误后点击 **开始生成**"
                )
                assistant_msg = Message(
                    conversation_id=conversation.id,
                    role="assistant",
                    content=summary
                )
                db.add(assistant_msg)
                await db.commit()

            yield {
                "event": "done",
                "data": json.dumps({"type": "done"})
            }
        except Exception as e:
            yield {
                "event": "error",
                "data": json.dumps({"type": "error", "message": str(e)})
            }

    return EventSourceResponse(event_generator())
