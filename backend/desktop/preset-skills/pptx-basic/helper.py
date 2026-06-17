from pptx import Presentation
from pptx.util import Pt

# [(标题, [要点...])]，第一项作封面（要点为空）
SLIDES = [
    ("示例标题", []),
    ("第一节", ["要点一", "要点二"]),
]

prs = Presentation()
for i, (title, bullets) in enumerate(SLIDES):
    layout = prs.slide_layouts[0 if i == 0 else 1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if bullets and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph(); p.text = b; p.font.size = Pt(18)
prs.save("output.pptx")
print("saved output.pptx")
