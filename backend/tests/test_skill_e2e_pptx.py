import pytest
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, create_async_engine, async_sessionmaker
from sqlalchemy.pool import StaticPool

from app.database import Base
from app.models import AIChatSession, AIChatArtifact
from app.ai_chat.tools import execute_tool


@pytest.mark.asyncio
async def test_pptx_skill_end_to_end(tmp_path, monkeypatch):
    pytest.importorskip("pptx")
    # 准备一个真 skill
    sk = tmp_path / "skills" / "platform" / "e2e-ppt"
    sk.mkdir(parents=True)
    sk.joinpath("SKILL.md").write_text("---\nname: e2e-ppt\ndescription: 测试PPT\n---\n跑 gen.py\n", encoding="utf-8")
    sk.joinpath("gen.py").write_text(
        "from pptx import Presentation\n"
        "p=Presentation(); s=p.slides.add_slide(p.slide_layouts[0]); s.shapes.title.text='Hi'\n"
        "p.save('output.pptx'); print('ok')\n", encoding="utf-8")
    monkeypatch.setenv("RUIJING_SKILLS_DIR", str(tmp_path / "skills"))

    engine = create_async_engine("sqlite+aiosqlite://", connect_args={"check_same_thread": False}, poolclass=StaticPool)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    Session = async_sessionmaker(engine, expire_on_commit=False, class_=AsyncSession)
    db = Session()
    ws = tmp_path / "ws"; ws.mkdir()
    s = AIChatSession(tenant_id=1, user_id=1, workspace_dir=str(ws))
    db.add(s); await db.commit(); await db.refresh(s)

    # 1) use_skill 展开 + 拷文件
    r1 = await execute_tool("use_skill", {"name": "e2e-ppt"}, s, db)
    assert "gen.py" in r1
    # 2) run_python 跑（dev 态用 venv python，python-pptx 可用）
    code = "import runpy; runpy.run_path('skill_e2e-ppt/gen.py', run_name='__main__')"
    r2 = await execute_tool("run_python", {"code": code}, s, db)
    assert "ok" in r2
    assert (ws / "output.pptx").is_file()
    # 3) 登记二进制产物
    r3 = await execute_tool("save_binary_artifact", {"source_path": "output.pptx", "filename": "测试.pptx"}, s, db)
    assert "测试.pptx" in r3
    art = (await db.execute(select(AIChatArtifact).where(AIChatArtifact.session_id == s.id))).scalar_one()
    assert art.storage == "file" and art.format == "pptx" and art.size_bytes > 0
